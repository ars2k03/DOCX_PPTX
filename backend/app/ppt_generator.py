import copy
import re
import time
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE


class PPTGenerationError(Exception):
    pass


def safe_text(value) -> str:
    if value is None:
        return ""
    return str(value).strip()


def duplicate_slide(prs: Presentation, source_slide):
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    for shape in list(new_slide.shapes):
        shape_element = shape.element
        shape_element.getparent().remove(shape_element)

    for shape in source_slide.shapes:
        copied_shape = copy.deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(
            copied_shape,
            "p:extLst",
        )

    return new_slide


def delete_slide(prs: Presentation, slide) -> None:
    slide_list = list(prs.slides)

    if slide not in slide_list:
        return

    slide_index = slide_list.index(slide)
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[slide_index]
    rel_id = slide_id.rId

    prs.part.drop_rel(rel_id)
    slide_id_list.remove(slide_id)


def remove_shape(shape) -> None:
    shape_element = shape.element
    shape_element.getparent().remove(shape_element)


def get_shape_text(shape) -> str:
    texts: List[str] = []

    if hasattr(shape, "shapes"):
        for child_shape in shape.shapes:
            texts.append(get_shape_text(child_shape))
        return " ".join(texts)

    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            texts.append(paragraph.text or "")

    return " ".join(texts)


def shape_contains_any_text(shape, needles: set[str]) -> bool:
    text = get_shape_text(shape)
    return any(needle in text for needle in needles)


def shape_vertical_bounds(shape) -> Optional[Tuple[int, int]]:
    if not hasattr(shape, "top") or not hasattr(shape, "height"):
        return None

    top = int(shape.top)
    bottom = int(shape.top + shape.height)

    return top, bottom


def row_overlaps(shape, row_top: int, row_bottom: int, margin: int = 50000) -> bool:
    bounds = shape_vertical_bounds(shape)

    if bounds is None:
        return False

    shape_top, shape_bottom = bounds

    expanded_top = row_top - margin
    expanded_bottom = row_bottom + margin

    return shape_bottom >= expanded_top and shape_top <= expanded_bottom


def find_optional_e_row(slide) -> Optional[Tuple[int, int]]:
    e_placeholders = {
        "{{OPTION_E}}",
        "{{OPTION_E_KEY}}",
        "{{OPTION_E_MARK}}",
        "{{OPTION_E_ROW}}",
    }

    matching_bounds: List[Tuple[int, int]] = []

    for shape in slide.shapes:
        if shape_contains_any_text(shape, e_placeholders):
            bounds = shape_vertical_bounds(shape)
            if bounds:
                matching_bounds.append(bounds)

    if not matching_bounds:
        return None

    row_top = min(top for top, _ in matching_bounds)
    row_bottom = max(bottom for _, bottom in matching_bounds)

    return row_top, row_bottom


def remove_optional_e_shapes(slide, mcq: Dict) -> None:
    option_e = safe_text(mcq.get("option_e", ""))

    if option_e:
        return

    e_row = find_optional_e_row(slide)

    if not e_row:
        return

    e_top, e_bottom = e_row

    for shape in list(slide.shapes):
        if row_overlaps(shape, e_top, e_bottom):
            remove_shape(shape)


def enable_text_fit(shape) -> None:
    if hasattr(shape, "shapes"):
        for child_shape in shape.shapes:
            enable_text_fit(child_shape)
        return

    if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
        return

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0


def set_text_font_size(shape, font_size: int) -> None:
    if hasattr(shape, "shapes"):
        for child_shape in shape.shapes:
            set_text_font_size(child_shape, font_size)
        return

    if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
        return

    enable_text_fit(shape)

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)


def estimate_required_lines(text: str, chars_per_line: int) -> int:
    text = safe_text(text)

    if not text:
        return 1

    words = text.split()
    lines = 1
    current_len = 0

    for word in words:
        word_len = len(word)

        if current_len + word_len + 1 <= chars_per_line:
            current_len += word_len + 1
        else:
            lines += 1
            current_len = word_len

    return max(lines, 1)


def estimate_chars_per_line(shape, font_size: int) -> int:
    if not hasattr(shape, "width"):
        return 12

    width_pt = shape.width / 12700
    avg_char_width_pt = font_size * 0.62

    if avg_char_width_pt <= 0:
        return 12

    return max(int(width_pt / avg_char_width_pt), 4)


def text_fits_shape(
        shape,
        text: str,
        font_size: int,
        line_height_multiplier: float = 1.18,
) -> bool:
    if not hasattr(shape, "height") or not hasattr(shape, "width"):
        return True

    height_pt = shape.height / 12700
    chars_per_line = estimate_chars_per_line(shape, font_size)
    lines = estimate_required_lines(text, chars_per_line)
    estimated_height_pt = lines * font_size * line_height_multiplier

    return estimated_height_pt <= height_pt


def fit_text_inside_existing_box(
        shape,
        text: str,
        default_font_size: int,
        min_font_size: int,
        line_height_multiplier: float = 1.18,
) -> int:
    if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
        return default_font_size

    enable_text_fit(shape)

    selected_font_size = min_font_size

    for font_size in range(default_font_size, min_font_size - 1, -2):
        if text_fits_shape(
                shape=shape,
                text=text,
                font_size=font_size,
                line_height_multiplier=line_height_multiplier,
        ):
            selected_font_size = font_size
            break

    set_text_font_size(shape, selected_font_size)
    return selected_font_size


def find_shapes_by_placeholder(slide, placeholder: str):
    matched = []

    def walk(shape):
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                walk(child)
            return

        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            if placeholder in get_shape_text(shape):
                matched.append(shape)

    for shape in slide.shapes:
        walk(shape)

    return matched


def get_answer_text(answer: str, mcq: Dict) -> str:
    answer = safe_text(answer).upper()

    option_map = {
        "A": safe_text(mcq.get("option_a", "")),
        "B": safe_text(mcq.get("option_b", "")),
        "C": safe_text(mcq.get("option_c", "")),
        "D": safe_text(mcq.get("option_d", "")),
        "E": safe_text(mcq.get("option_e", "")),
    }

    return option_map.get(answer, "")


def get_valid_answer_key(mcq: Dict) -> str:
    answer = safe_text(mcq.get("answer", "")).upper()

    option_map = {
        "A": safe_text(mcq.get("option_a", "")),
        "B": safe_text(mcq.get("option_b", "")),
        "C": safe_text(mcq.get("option_c", "")),
        "D": safe_text(mcq.get("option_d", "")),
        "E": safe_text(mcq.get("option_e", "")),
    }

    if answer not in option_map:
        return ""

    if not option_map.get(answer):
        return ""

    return answer


def get_question_display_text(mcq: Dict) -> str:
    question = safe_text(mcq.get("question", ""))
    source = safe_text(mcq.get("question_source", ""))

    if source:
        return f"{question} [{source}]"

    return question


def apply_auto_fit_fonts(slide, mcq: Dict) -> None:
    question_text = safe_text(mcq.get("question", ""))

    for shape in find_shapes_by_placeholder(slide, "{{QUESTION}}"):
        fit_text_inside_existing_box(
            shape=shape,
            text=question_text,
            default_font_size=96,
            min_font_size=20,
            line_height_multiplier=1.15,
        )

    option_map = {
        "{{OPTION_A}}": safe_text(mcq.get("option_a", "")),
        "{{OPTION_B}}": safe_text(mcq.get("option_b", "")),
        "{{OPTION_C}}": safe_text(mcq.get("option_c", "")),
        "{{OPTION_D}}": safe_text(mcq.get("option_d", "")),
        "{{OPTION_E}}": safe_text(mcq.get("option_e", "")),
    }

    for placeholder, option_text in option_map.items():
        for shape in find_shapes_by_placeholder(slide, placeholder):
            fit_text_inside_existing_box(
                shape=shape,
                text=option_text,
                default_font_size=88,
                min_font_size=16,
                line_height_multiplier=1.18,
            )

    answer = get_valid_answer_key(mcq)
    answer_text = get_answer_text(answer, mcq)

    for shape in find_shapes_by_placeholder(slide, "{{ANSWER_TEXT}}"):
        fit_text_inside_existing_box(
            shape=shape,
            text=answer_text,
            default_font_size=88,
            min_font_size=16,
            line_height_multiplier=1.18,
        )

    for shape in find_shapes_by_placeholder(slide, "{{ANSWER_KEY}}"):
        fit_text_inside_existing_box(
            shape=shape,
            text=answer,
            default_font_size=88,
            min_font_size=24,
            line_height_multiplier=1.05,
        )


def clean_empty_source_brackets(text: str) -> str:
    text = text.replace("[]", "")
    text = text.replace("[ ]", "")
    text = re.sub(r"\s+\]", "]", text)
    text = re.sub(r"\[\s+", "[", text)
    text = re.sub(r"\s{2,}", " ", text)
    return text.strip()


def replace_text_in_text_frame(text_frame, replacements: Dict[str, str]) -> None:
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0

    for paragraph in text_frame.paragraphs:
        runs = paragraph.runs

        if not runs:
            continue

        original_text = "".join(run.text for run in runs)

        if not original_text:
            continue

        new_text = original_text

        for key, value in replacements.items():
            new_text = new_text.replace(key, safe_text(value))

        new_text = clean_empty_source_brackets(new_text)

        if new_text == original_text:
            continue

        runs[0].text = new_text

        for run in runs[1:]:
            run.text = ""


def replace_text_in_shape(shape, replacements: Dict[str, str]) -> None:
    if hasattr(shape, "shapes"):
        for child_shape in shape.shapes:
            replace_text_in_shape(child_shape, replacements)
        return

    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        replace_text_in_text_frame(shape.text_frame, replacements)


def build_replacements(mcq: Dict, show_answer: bool) -> Dict[str, str]:
    option_a = safe_text(mcq.get("option_a", ""))
    option_b = safe_text(mcq.get("option_b", ""))
    option_c = safe_text(mcq.get("option_c", ""))
    option_d = safe_text(mcq.get("option_d", ""))
    option_e = safe_text(mcq.get("option_e", ""))

    answer = get_valid_answer_key(mcq)
    answer_text = get_answer_text(answer, mcq)

    source = safe_text(mcq.get("question_source", ""))
    question = safe_text(mcq.get("question", ""))
    question_full = get_question_display_text(mcq)

    return {
        "{{QUESTION_SOURCE}}": source,
        "[{{QUESTION_SOURCE}}]": f"[{source}]" if source else "",

        # Keep {{QUESTION}} source-free because most templates already place
        # [{{QUESTION_SOURCE}}] beside/after it. This prevents double source like:
        # Question [DU...] [DU...]
        "{{QUESTION}}": question,
        "{{QUESTION_ONLY}}": question,
        "{{QUESTION_FULL}}": question_full,

        "{{OPTION_A}}": option_a,
        "{{OPTION_B}}": option_b,
        "{{OPTION_C}}": option_c,
        "{{OPTION_D}}": option_d,
        "{{OPTION_E}}": option_e,

        "{{OPTION_A_KEY}}": "A",
        "{{OPTION_B_KEY}}": "B",
        "{{OPTION_C_KEY}}": "C",
        "{{OPTION_D_KEY}}": "D",
        "{{OPTION_E_KEY}}": "E" if option_e else "",

        "{{ANSWER_KEY}}": answer if show_answer else "",
        "{{ANSWER_TEXT}}": answer_text if show_answer else "",

        "{{OPTION_A_MARK}}": "✓ " if show_answer and answer == "A" else "",
        "{{OPTION_B_MARK}}": "✓ " if show_answer and answer == "B" else "",
        "{{OPTION_C_MARK}}": "✓ " if show_answer and answer == "C" else "",
        "{{OPTION_D_MARK}}": "✓ " if show_answer and answer == "D" else "",
        "{{OPTION_E_MARK}}": "✓ " if show_answer and answer == "E" else "",

        "{{OPTION_E_ROW}}": "",
    }


def replace_slide_placeholders(slide, mcq: Dict, show_answer: bool) -> None:
    remove_optional_e_shapes(slide, mcq)
    apply_auto_fit_fonts(slide, mcq)

    replacements = build_replacements(
        mcq=mcq,
        show_answer=show_answer,
    )

    for shape in slide.shapes:
        replace_text_in_shape(shape, replacements)


def generate_pptx(
        mcqs: List[Dict],
        template_dir: str | Path = "templates",
        output_dir: str | Path = "outputs",
) -> Path:
    if not mcqs:
        raise PPTGenerationError("No MCQs provided.")

    template_dir = Path(template_dir)
    output_dir = Path(output_dir)
    template_path = template_dir / "template.pptx"

    if not template_path.exists():
        raise PPTGenerationError(f"Template file not found: {template_path}")

    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(template_path))
    template_slides = list(prs.slides)

    if len(template_slides) < 2:
        raise PPTGenerationError(
            "template.pptx must contain 2 slides: "
            "slide 1 question-only, slide 2 answer-box slide."
        )

    question_template_slide = template_slides[0]
    answer_template_slide = template_slides[1]

    for mcq in mcqs:
        question_slide = duplicate_slide(prs, question_template_slide)
        replace_slide_placeholders(
            slide=question_slide,
            mcq=mcq,
            show_answer=False,
        )

        answer_slide = duplicate_slide(prs, answer_template_slide)
        replace_slide_placeholders(
            slide=answer_slide,
            mcq=mcq,
            show_answer=True,
        )

    for slide in template_slides:
        delete_slide(prs, slide)

    timestamp = int(time.time())
    total_slides = len(mcqs) * 2

    output_name = f"MCQ_{len(mcqs)}Q_{total_slides}Slides_{timestamp}.pptx"
    output_path = output_dir / output_name

    prs.save(output_path)

    return output_path